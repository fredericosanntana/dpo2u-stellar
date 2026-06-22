from pathlib import Path
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path('/root/dpo2u-stellar')
OUT = ROOT / 'artifacts' / 'pulso'
OUT.mkdir(parents=True, exist_ok=True)

PDF_PATH = OUT / 'Pulso-One-Pager.pdf'
PPTX_PATH = OUT / 'Pulso-Deck.pptx'

# Palette
NAVY = RGBColor(13, 27, 42)
BLUE = RGBColor(27, 75, 132)
TEAL = RGBColor(36, 123, 160)
MINT = RGBColor(109, 203, 176)
SAND = RGBColor(245, 245, 240)
RED = RGBColor(180, 56, 72)
WHITE = RGBColor(255, 255, 255)
CHAR = RGBColor(33, 37, 41)
GREY = RGBColor(106, 117, 127)
LIGHT = RGBColor(236, 242, 248)


def add_textbox(slide, x, y, w, h, text, size=20, color=CHAR, bold=False, font='Aptos',
                align=PP_ALIGN.LEFT, margin=0.08, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_circle(slide, x, y, d, fill, text='', text_color=WHITE, size=18):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    if text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(size)
        f.bold = True
        f.color.rgb = text_color
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def one_pager():
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(PDF_PATH), pagesize=A4,
                            leftMargin=1.2*cm, rightMargin=1.2*cm,
                            topMargin=1.1*cm, bottomMargin=1.1*cm)
    title = ParagraphStyle('title', parent=styles['Title'], fontName='Helvetica-Bold',
                           fontSize=23, leading=27, textColor=colors.HexColor('#0D1B2A'), alignment=TA_LEFT)
    subtitle = ParagraphStyle('subtitle', parent=styles['BodyText'], fontName='Helvetica',
                              fontSize=10, leading=13, textColor=colors.HexColor('#4A5568'))
    h = ParagraphStyle('h', parent=styles['Heading2'], fontName='Helvetica-Bold',
                       fontSize=11.5, leading=14, textColor=colors.HexColor('#1B4B84'), spaceAfter=4)
    body = ParagraphStyle('body', parent=styles['BodyText'], fontName='Helvetica',
                          fontSize=9.4, leading=12.2, textColor=colors.HexColor('#1F2937'))
    body_small = ParagraphStyle('body_small', parent=body, fontSize=8.8, leading=11.2)
    quote = ParagraphStyle('quote', parent=body, fontName='Helvetica-Bold', fontSize=11,
                           leading=14, textColor=colors.HexColor('#0D1B2A'), leftIndent=8)

    story = []
    story.append(Paragraph('Pulso — One Pager Executivo', title))
    story.append(Paragraph('DPO2U | Stellar | junho de 2026', subtitle))
    story.append(Spacer(1, 0.25*cm))
    story.append(Paragraph('Tese central', h))
    story.append(Paragraph('A DPO2U transforma decisões canônicas de compliance em consequências operacionais verificáveis on-chain. No Pulso, a prova principal não é um dashboard: é uma lane institucional em Stellar onde o registry decide, a admissão acontece e a revogação produz bloqueio verificável.', quote))
    story.append(Spacer(1, 0.22*cm))

    data = [
        [Paragraph('<b>O que já está provado</b>', h), Paragraph('<b>Por que isso importa</b>', h), Paragraph('<b>Boundary honesto</b>', h)],
        [Paragraph('• Registry vivo em testnet<br/>• Admissão na lane ASP/SPP<br/>• Revogação canônica on-chain<br/>• Blocked-lane real em instância própria<br/>• Worker idempotente em rerun', body),
         Paragraph('• Compliance deixa de ser PDF<br/>• A decisão canônica muda o comportamento da infraestrutura<br/>• A lane é replayável e verificável<br/>• A primitive já sustenta uma narrativa institucional clara', body),
         Paragraph('• A instância externa é auditável, não operada por nós<br/>• Mutação externa ainda depende de governança/admin authority<br/>• DeFindex entra como supporting evidence, não como narrativa principal', body)]
    ]
    t = Table(data, colWidths=[6.0*cm, 6.0*cm, 6.0*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#EAF1F8')),
        ('BOX', (0,0), (-1,-1), 0.8, colors.HexColor('#C7D3E0')),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#D7E0EA')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 8),
        ('RIGHTPADDING', (0,0), (-1,-1), 8),
        ('TOPPADDING', (0,0), (-1,-1), 7),
        ('BOTTOMPADDING', (0,0), (-1,-1), 7),
        ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#FBFCFE')),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.3*cm))

    story.append(Paragraph('As 3 vertentes', h))
    data2 = [
        [Paragraph('<b>1. Pulso hackathon</b>', body_small), Paragraph('Demo principal: registry-backed admission + revocation consequence + blocked-lane enforcement na lane própria B-first.', body_small)],
        [Paragraph('<b>2. Hackathon ZK Stellar</b>', body_small), Paragraph('Extensão da mesma primitive para elegibilidade/admissão com privacidade. Não deve competir com a narrativa central do Pulso.', body_small)],
        [Paragraph('<b>3. GTM da solução</b>', body_small), Paragraph('Empacotar a tese como “compliance as infrastructure”: lane verificável, governance boundary honesto e encaixe institucional claro para parceiros/exchanges.', body_small)],
    ]
    t2 = Table(data2, colWidths=[4.2*cm, 13.8*cm])
    t2.setStyle(TableStyle([
        ('BOX', (0,0), (-1,-1), 0.8, colors.HexColor('#C7D3E0')),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#D7E0EA')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 8),
        ('RIGHTPADDING', (0,0), (-1,-1), 8),
        ('TOPPADDING', (0,0), (-1,-1), 7),
        ('BOTTOMPADDING', (0,0), (-1,-1), 7),
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#F9FBFD')),
    ]))
    story.append(t2)
    story.append(Spacer(1, 0.25*cm))

    story.append(Paragraph('Mensagem para parceiro/judge', h))
    story.append(Paragraph('“Não estamos vendendo compliance em dashboard. Estamos mostrando que uma decisão canônica de compliance já consegue produzir consequência operacional verificável em Stellar, com boundary honesto e sem overengineering narrativo.”', body))
    story.append(Spacer(1, 0.18*cm))
    story.append(Paragraph('Próximo passo recomendado', h))
    story.append(Paragraph('Usar Pulso como wedge demonstrativo, ZK como aprofundamento da mesma primitive e GTM como embalagem institucional da tese — sem misturar os três em uma única demo.', body))
    doc.build(story)


def style_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(8.0), Inches(1.0), title, size=28, color=WHITE, bold=True, font='Aptos Display')
    if subtitle:
        add_textbox(slide, Inches(0.72), Inches(1.35), Inches(7.5), Inches(0.5), subtitle, size=12.5, color=RGBColor(210,220,235))


def add_footer(slide, text='DPO2U • Pulso • junho 2026'):
    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(2.8), Inches(0.25), text, size=9, color=GREY)


def deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    add_rect(s, Inches(8.7), Inches(0.8), Inches(3.7), Inches(5.9), TEAL, radius=True)
    add_rect(s, Inches(9.0), Inches(1.15), Inches(3.1), Inches(1.0), MINT, radius=True)
    add_rect(s, Inches(9.0), Inches(2.45), Inches(3.1), Inches(1.0), BLUE, radius=True)
    add_rect(s, Inches(9.0), Inches(3.75), Inches(3.1), Inches(1.0), RGBColor(27, 75, 132), radius=True)
    add_textbox(s, Inches(9.18), Inches(1.42), Inches(2.7), Inches(0.4), 'Registry', size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(9.18), Inches(2.72), Inches(2.7), Inches(0.4), 'Admission', size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(9.18), Inches(4.02), Inches(2.7), Inches(0.4), 'Blocked lane', size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    style_title(s, 'Pulso', 'Compliance como consequência operacional verificável em Stellar')
    add_textbox(s, Inches(0.75), Inches(2.15), Inches(6.4), Inches(1.8), 'A DPO2U transforma decisões canônicas de compliance em admissão, revogação e bloqueio verificáveis on-chain.', size=22, color=WHITE, bold=False)
    add_textbox(s, Inches(0.76), Inches(4.65), Inches(4.5), Inches(0.5), 'Pulso hackathon • ZK track • GTM institucional', size=13, color=RGBColor(190,205,220))

    # Slide 2
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.45), Inches(5.0), Inches(0.5), 'A tese em uma frase', size=26, color=NAVY, bold=True, font='Aptos Display')
    add_rect(s, Inches(0.7), Inches(1.15), Inches(12.0), Inches(1.2), LIGHT, radius=True)
    add_textbox(s, Inches(0.95), Inches(1.42), Inches(11.3), Inches(0.7), 'Registry-backed admission + revocation consequence + blocked-lane enforcement na lane própria B-first.', size=21, color=NAVY, bold=True)
    cards = [
        ('Problema', 'Compliance tradicional para no parecer e não muda comportamento operacional.'),
        ('Tese', 'A decisão canônica precisa mover a infraestrutura, não só a documentação.'),
        ('Demonstração', 'Pulso prova isso em Stellar com lane replayável, verificável e soberana.')
    ]
    xs = [0.7, 4.45, 8.2]
    for i, (k, v) in enumerate(cards):
        add_rect(s, Inches(xs[i]), Inches(2.8), Inches(3.35), Inches(2.6), SAND, line=RGBColor(220,226,232), radius=True)
        add_circle(s, Inches(xs[i]+0.22), Inches(3.0), Inches(0.45), TEAL, str(i+1), size=16)
        add_textbox(s, Inches(xs[i]+0.8), Inches(3.0), Inches(2.2), Inches(0.35), k, size=18, color=BLUE, bold=True)
        add_textbox(s, Inches(xs[i]+0.22), Inches(3.55), Inches(2.85), Inches(1.4), v, size=15, color=CHAR)
    add_footer(s)

    # Slide 3
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.45), Inches(6.2), Inches(0.5), 'O que já está provado no repo', size=26, color=NAVY, bold=True, font='Aptos Display')
    items = [
        ('Registry vivo', 'Contract ativo em testnet com register, verify e revoke reais.'),
        ('Admissão na lane', 'Decisão do registry já vira artefato operacional de admissão ASP/SPP.'),
        ('Blocked-lane real', 'Instância própria do asp-non-membership já bloqueia e reverte on-chain.'),
        ('Worker idempotente', 'No-op quando ativo, bloqueio quando revogado, rerun sem duplicação.'),
    ]
    for i, (k, v) in enumerate(items):
        top = 1.4 + i*1.28
        add_circle(s, Inches(0.85), Inches(top), Inches(0.42), NAVY, str(i+1), size=16)
        add_textbox(s, Inches(1.45), Inches(top-0.02), Inches(3.0), Inches(0.3), k, size=18, color=BLUE, bold=True)
        add_textbox(s, Inches(1.45), Inches(top+0.33), Inches(5.1), Inches(0.55), v, size=14.5, color=CHAR)
        add_rect(s, Inches(7.2), Inches(top-0.05), Inches(5.2), Inches(0.9), [TEAL, BLUE, MINT, RED][i], radius=True)
        value = [
            'CAUD…WTYP',
            'decision → admission → record',
            'insert / delete / verify',
            'no-op / block / rerun-safe'
        ][i]
        add_textbox(s, Inches(7.45), Inches(top+0.19), Inches(4.7), Inches(0.4), value, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)

    # Slide 4
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    add_textbox(s, Inches(0.7), Inches(0.45), Inches(6.5), Inches(0.5), 'A cadeia operacional do Pulso', size=26, color=WHITE, bold=True, font='Aptos Display')
    steps = [
        ('1', 'Registry', 'A attestation nasce como verdade canônica.'),
        ('2', 'Admission', 'A decisão vira admissão operacional na lane.'),
        ('3', 'Revocation', 'A prova deixa de valer e a lane muda de estado.'),
        ('4', 'Blocked lane', 'A consequência de bloqueio se torna verificável on-chain.'),
    ]
    xs = [0.8, 3.5, 6.2, 8.9]
    colors_ = [MINT, TEAL, RGBColor(219, 142, 68), RED]
    for i, (n, t, d) in enumerate(steps):
        add_rect(s, Inches(xs[i]), Inches(2.0), Inches(2.25), Inches(2.9), WHITE if i % 2 == 0 else LIGHT, line=WHITE, radius=True)
        add_circle(s, Inches(xs[i]+0.82), Inches(2.2), Inches(0.55), colors_[i], n, size=20)
        add_textbox(s, Inches(xs[i]+0.18), Inches(2.95), Inches(1.9), Inches(0.55), t, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(xs[i]+0.2), Inches(3.52), Inches(1.85), Inches(0.95), d, size=13, color=CHAR, align=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(s, Inches(xs[i]+2.28), Inches(3.1), Inches(0.45), Inches(0.4), '→', size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.85), Inches(5.7), Inches(11.6), Inches(0.6), 'Não é compliance em PDF. É uma decisão canônica produzindo consequência operacional em Stellar.', size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Slide 5
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.45), Inches(8.4), Inches(0.8), 'Boundary honesto e anti-overengineering', size=22, color=NAVY, bold=True, font='Aptos Display')
    add_rect(s, Inches(0.7), Inches(1.55), Inches(5.7), Inches(4.65), LIGHT, line=RGBColor(210,220,230), radius=True)
    add_rect(s, Inches(6.9), Inches(1.55), Inches(5.7), Inches(4.65), SAND, line=RGBColor(220,226,232), radius=True)
    add_textbox(s, Inches(0.95), Inches(1.88), Inches(4.5), Inches(0.4), 'O que podemos afirmar', size=19, color=BLUE, bold=True)
    add_textbox(s, Inches(1.0), Inches(2.42), Inches(4.8), Inches(2.65), '• Operamos a lane própria B-first\n• Registry, admission, revoke e blocked-lane já foram provados\n• Worker idempotente já existe\n• A primitive já tem supporting evidence em DeFindex', size=15, color=CHAR)
    add_textbox(s, Inches(7.15), Inches(1.82), Inches(4.75), Inches(0.7), 'O que não devemos\noverclaimar', size=18, color=RED, bold=True)
    add_textbox(s, Inches(7.2), Inches(2.7), Inches(4.8), Inches(2.35), '• Não operamos a instância externa auditada\n• O gap remanescente é governança, não integração básica\n• O track ZK não é outro produto\n• DeFindex não é a narrativa principal do Pulso', size=15, color=CHAR)
    add_footer(s)

    # Slide 6
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.45), Inches(6.2), Inches(0.5), 'As 3 vertentes', size=26, color=NAVY, bold=True, font='Aptos Display')
    rows = [
        ('Pulso hackathon', 'Fechar demo principal', 'Registry-backed admission + revocation + blocked-lane na lane própria.'),
        ('Hackathon ZK Stellar', 'Estender a primitive', 'Privacidade/elegibilidade como aprofundamento da mesma tese, sem competir com o Pulso.'),
        ('GTM da solução', 'Empacotar para mercado', 'Vender “compliance as infrastructure” com evidência real e boundary honesto.'),
    ]
    for i, (a, b, c) in enumerate(rows):
        top = 1.5 + i*1.65
        add_rect(s, Inches(0.8), Inches(top), Inches(11.8), Inches(1.15), LIGHT if i != 1 else SAND, line=RGBColor(215,223,232), radius=True)
        add_rect(s, Inches(0.8), Inches(top), Inches(2.55), Inches(1.15), [TEAL, MINT, BLUE][i], line=[TEAL, MINT, BLUE][i], radius=True)
        add_textbox(s, Inches(1.02), Inches(top+0.28), Inches(2.0), Inches(0.35), a, size=18, color=WHITE, bold=True)
        add_textbox(s, Inches(3.65), Inches(top+0.16), Inches(2.4), Inches(0.3), b, size=17, color=NAVY, bold=True)
        add_textbox(s, Inches(6.15), Inches(top+0.14), Inches(5.95), Inches(0.6), c, size=14.5, color=CHAR)
    add_footer(s)

    # Slide 7
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    add_rect(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(3.9), RGBColor(18,38,58), radius=True)
    add_textbox(s, Inches(1.3), Inches(2.0), Inches(10.5), Inches(1.0), 'Mensagem final', size=24, color=MINT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.45), Inches(2.9), Inches(10.2), Inches(1.4), 'A DPO2U já consegue provar em Stellar que uma decisão canônica de compliance pode produzir consequência operacional verificável. Pulso é a wedge demo. ZK é a extensão. GTM é a embalagem institucional.', size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(2.5), Inches(5.9), Inches(8.2), Inches(0.5), 'Compliance as infrastructure — não compliance como dashboard.', size=16, color=RGBColor(205,220,235), bold=False, align=PP_ALIGN.CENTER)

    prs.save(str(PPTX_PATH))


if __name__ == '__main__':
    one_pager()
    deck()
    print(PDF_PATH)
    print(PPTX_PATH)
